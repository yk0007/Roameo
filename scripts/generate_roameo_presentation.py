from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTDIR = Path("presentation")
OUTFILE = OUTDIR / "Roameo-Final-Year-Presentation.pptx"


slides = [
    {
        "title": "Roameo",
        "subtitle": "An AI-Powered Travel Planning Workspace with Canonical Session State",
        "bullets": [
            "Final Year Project Presentation",
            "Text-first planning workspace with chat, itinerary, and map",
            "Built from the actual Roameo codebase architecture"
        ]
    },
    {
        "title": "Problem Statement",
        "bullets": [
            "Trip planning is fragmented across search, maps, notes, and booking apps",
            "Typical AI chatbots suggest places but do not create a persistent editable plan",
            "Users face inconsistency between recommendations, routes, and saved items"
        ]
    },
    {
        "title": "Proposed Solution",
        "bullets": [
            "Roameo is a text-first AI travel planning workspace",
            "Natural language request becomes a structured itinerary and live map view",
            "One canonical session snapshot keeps all product surfaces synchronized"
        ]
    },
    {
        "title": "Project Objectives",
        "bullets": [
            "Support multi-turn conversational trip planning and refinement",
            "Maintain one source of truth for chat, itinerary, map, and saved POIs",
            "Use AI for semantics and synthesis, deterministic code for state mutation",
            "Ground outputs with real-world travel APIs and stream updates live"
        ]
    },
    {
        "title": "Core Innovation",
        "subtitle": "Canonical Session Snapshot",
        "bullets": [
            "Single SessionSnapshot drives messages, plan, POI catalog, saved items, memory, and traces",
            "Prevents chat/map/itinerary mismatches common in AI apps",
            "Makes follow-up turns reliable because context persists structurally"
        ]
    },
    {
        "title": "High-Level Architecture",
        "bullets": [
            "Frontend: Next.js 16 + React 19 workspace UI",
            "Backend: Express + TypeScript runtime and APIs",
            "Shared contracts: Zod schemas + TypeScript types",
            "Persistence: Supabase with in-memory fallback for development"
        ]
    },
    {
        "title": "Technology Stack",
        "bullets": [
            "Frontend: Next.js, React, Zustand, React Query",
            "Backend: Express, TypeScript, Zod, Supabase",
            "AI Providers: Gemini and OpenAI",
            "Travel Integrations: Google Places, Geocoding, Directions, Open-Meteo, Tavily, Nager.Date"
        ]
    },
    {
        "title": "Product Interface",
        "bullets": [
            "Chat panel on the left for natural language planning",
            "Map or itinerary panel on the right for structured visualization",
            "Top navigation for destination, dates, travelers, and budget",
            "Designed as a workspace, not just a chatbot"
        ]
    },
    {
        "title": "End-to-End Request Flow",
        "bullets": [
            "Save user message and mark planning state as running",
            "Resolve intent and context with fast-path + semantic routing",
            "Run discovery and enrichment when needed",
            "Synthesize, validate, narrate, persist, and stream the final state"
        ]
    },
    {
        "title": "Agentic Runtime Design",
        "bullets": [
            "Router-first architecture with deterministic orchestration",
            "LLM used for intent resolution, structured planning, and narrative generation",
            'Everything else stays deterministic: provider selection, tool execution, state writes, stream events'
        ]
    },
    {
        "title": "Specialized Sub-Agents",
        "bullets": [
            "Semantic router for travel intent and destination context",
            "Destination research and discovery handling",
            "Structured itinerary synthesis",
            "Feasibility critic and transit advisor for practical refinement"
        ]
    },
    {
        "title": "Internal Tool and Mutation Layer",
        "bullets": [
            "Safe tools for reading snapshot, editing itinerary, and updating trip metadata",
            "Canonical plan mutations: add/remove POI, move activity, regenerate day, rebalance trip",
            "Centralized writes keep session state consistent and auditable"
        ]
    },
    {
        "title": "Real-World Data Grounding",
        "bullets": [
            "Google Places for POIs such as stays, restaurants, and attractions",
            "Google Geocoding and Directions for location and route support",
            "Open-Meteo, holidays, and Tavily add weather, calendar, and editorial context",
            "Avoids fabricated travel recommendations"
        ]
    },
    {
        "title": "State Management Rules",
        "bullets": [
            "Explicit new trip requests replace stale active-trip context",
            "Multi-city trips preserve the full destination set",
            "Discovery expands the POI catalog without corrupting the itinerary",
            "Map routes come only from itinerary-linked POIs"
        ]
    },
    {
        "title": "Frontend as a Thin Consumer",
        "bullets": [
            "Loads session snapshot and subscribes to SSE updates",
            "Hydrates Zustand state and derives map/itinerary projections",
            "Presentation stays on the frontend; business logic stays on the backend"
        ]
    },
    {
        "title": "Live Planning Experience",
        "bullets": [
            "Server-Sent Events stream session.snapshot, plan.updated, trace.updated, and turn completion",
            "Users see progress live instead of waiting for one final static answer",
            "Improves transparency and perceived responsiveness"
        ]
    },
    {
        "title": "Persistence Strategy",
        "bullets": [
            "Stores sessions, messages, plan snapshots, POI catalogs, saved POIs, and traces",
            "Supabase-backed persistence for authenticated usage",
            "In-memory fallback supports local development and testing"
        ]
    },
    {
        "title": "Key Strengths",
        "bullets": [
            "Single-source architecture",
            "Hybrid AI + deterministic runtime",
            "Multi-turn contextual planning",
            "Real-world data grounding",
            "Live synchronized user experience"
        ]
    },
    {
        "title": "Limitations",
        "bullets": [
            "Output quality depends on external API availability and latency",
            "Travel data freshness varies by source",
            "Personalization can be improved with longer-term user modeling",
            "Deeper route and price optimization remain future work"
        ]
    },
    {
        "title": "Future Enhancements",
        "bullets": [
            "Flight and train integration",
            "Cost-aware optimization",
            "Collaborative group trip planning",
            "Multilingual support and booking assistance"
        ]
    },
    {
        "title": "Conclusion",
        "bullets": [
            "Roameo combines natural language interaction, tool-grounded planning, and canonical state management",
            "The main contribution is a reliable travel planning workspace, not just itinerary generation",
            "Chat, itinerary, map, and saved places remain synchronized across the full trip-planning flow"
        ]
    },
    {
        "title": "Demo / Q&A",
        "bullets": [
            "Show trip creation from a natural language prompt",
            "Show streaming plan generation and refinement",
            "Show map + itinerary consistency inside one live session"
        ]
    }
]


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 246, 241)


def add_title_band(slide, title, subtitle=None):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.3)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(17, 48, 59)
    band.line.color.rgb = RGBColor(17, 48, 59)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.28), Inches(8.8), Inches(0.45))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.76), Inches(10.5), Inches(0.25))
        p2 = sub_box.text_frame.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = "Aptos"
        run2.font.size = Pt(12)
        run2.font.color.rgb = RGBColor(217, 229, 234)


def add_bullets(slide, bullets):
    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.75), Inches(8.6), Inches(4.9))
    tf = body.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(12)
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(22)
            run.font.color.rgb = RGBColor(29, 39, 44)


def add_side_panel(slide, header, lines):
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.85), Inches(1.65), Inches(2.75), Inches(4.85)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(222, 232, 226)
    panel.line.color.rgb = RGBColor(164, 186, 176)

    hbox = slide.shapes.add_textbox(Inches(10.1), Inches(1.95), Inches(2.2), Inches(0.4))
    hp = hbox.text_frame.paragraphs[0]
    hr = hp.add_run()
    hr.text = header
    hr.font.name = "Aptos"
    hr.font.size = Pt(16)
    hr.font.bold = True
    hr.font.color.rgb = RGBColor(17, 48, 59)

    lbox = slide.shapes.add_textbox(Inches(10.1), Inches(2.45), Inches(2.1), Inches(3.65))
    ltf = lbox.text_frame
    ltf.word_wrap = True
    for idx, line in enumerate(lines):
        p = ltf.paragraphs[0] if idx == 0 else ltf.add_paragraph()
        p.text = line
        p.space_after = Pt(10)
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(13)
            run.font.color.rgb = RGBColor(39, 59, 64)


def add_footer(slide, idx, total):
    footer = slide.shapes.add_textbox(Inches(11.6), Inches(7.0), Inches(1.1), Inches(0.3))
    p = footer.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{idx}/{total}"
    r.font.name = "Aptos"
    r.font.size = Pt(10)
    r.font.color.rgb = RGBColor(110, 119, 123)


def main():
    OUTDIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]

    total = len(slides)
    for idx, content in enumerate(slides, start=1):
        slide = prs.slides.add_slide(layout)
        set_background(slide)
        add_title_band(slide, content["title"], content.get("subtitle"))
        add_bullets(slide, content["bullets"])
        add_side_panel(
            slide,
            "Presenter Cue",
            [
                "Explain this slide in simple terms.",
                "Relate the point back to consistency, reliability, or live planning.",
                "Use the demo to reinforce this claim if needed."
            ],
        )
        add_footer(slide, idx, total)

    prs.save(OUTFILE)
    print(f"Created {OUTFILE}")


if __name__ == "__main__":
    main()
